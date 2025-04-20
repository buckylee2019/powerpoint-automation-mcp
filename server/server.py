from mcp.server.fastmcp import FastMCP
import os
import uuid
from typing import Dict, List, Optional, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
import tempfile

mcp = FastMCP("ppts")

USER_AGENT = "ppts-app/1.0"

class PPTAutomation:
    def __init__(self):
        self.presentations = {}  # Store presentation IDs and their objects
        
    def initialize(self):
        """Initialize the PowerPoint automation - no app instance needed with python-pptx"""
        return True
                
    def get_open_presentations(self):
        """Get all currently tracked presentations"""
        result = []
        for pres_id, pres_info in self.presentations.items():
            result.append({
                "id": pres_id,
                "name": os.path.basename(pres_info["path"]) if pres_info["path"] else "Untitled",
                "path": pres_info["path"],
                "slide_count": len(pres_info["presentation"].slides)
            })
        return result

# Create a global instance of our automation class
ppt_automation = PPTAutomation()

@mcp.tool()
def initialize_powerpoint() -> bool:
    """Initialize connection to PowerPoint."""
    return ppt_automation.initialize()

@mcp.tool()
def get_presentations() -> List[Dict[str, Any]]:
    """Get a list of all tracked PowerPoint presentations with their metadata."""
    return ppt_automation.get_open_presentations()

@mcp.tool()
def open_presentation(file_path: str) -> Dict[str, Any]:
    """
    Open a PowerPoint presentation from the specified path.
    
    Args:
        file_path: Full path to the PowerPoint file (.pptx)
        
    Returns:
        Dictionary with presentation ID and metadata
    """
    if not os.path.exists(file_path):
        return {"error": f"File not found: {file_path}"}
    
    try:
        pres = Presentation(file_path)
        pres_id = str(uuid.uuid4())
        
        ppt_automation.presentations[pres_id] = {
            "presentation": pres,
            "path": file_path
        }
        
        return {
            "id": pres_id,
            "name": os.path.basename(file_path),
            "path": file_path,
            "slide_count": len(pres.slides)
        }
    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def get_slides(presentation_id: str) -> List[Dict[str, Any]]:
    """
    Get a list of all slides in a presentation.
    
    Args:
        presentation_id: ID of the presentation
        
    Returns:
        List of slide metadata
    """
    if presentation_id not in ppt_automation.presentations:
        return {"error": "Presentation ID not found"}
    
    pres = ppt_automation.presentations[presentation_id]["presentation"]
    slides = []
    
    try:
        for i, slide in enumerate(pres.slides):
            slide_id = str(i)
            
            # Try to get slide title
            title = "Untitled Slide"
            for shape in slide.shapes:
                if shape.is_placeholder and shape.placeholder_format.type == 1:  # Title placeholder
                    if shape.has_text_frame:
                        title = shape.text
                        break
            
            slides.append({
                "id": slide_id,
                "index": i,
                "title": title,
                "shape_count": len(slide.shapes)
            })
        
        return slides
    except Exception as e:
        return {"error": f"Error getting slides: {str(e)}"}

@mcp.tool()
def get_slide_text(presentation_id: str, slide_index: int) -> Dict[str, Any]:
    """
    Get all text content in a slide.
    
    Args:
        presentation_id: ID of the presentation
        slide_index: Index of the slide (integer, 0-based)
        
    Returns:
        Dictionary containing text content organized by shape
    """
    try:
        if presentation_id not in ppt_automation.presentations:
            return {"error": f"Presentation ID not found: {presentation_id}"}
        
        pres = ppt_automation.presentations[presentation_id]["presentation"]
        
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {"error": f"Invalid slide index: {slide_index}. Valid range is 0-{len(pres.slides)-1}"}
        
        slide = pres.slides[slide_index]
        text_content = {}
        
        for shape_idx, shape in enumerate(slide.shapes):
            shape_id = str(shape_idx)
            
            if shape.has_text_frame:
                text = shape.text
                shape_name = f"Shape {shape_idx}"
                
                text_content[shape_id] = {
                    "shape_name": shape_name,
                    "text": text
                }
        
        return {
            "slide_index": slide_index,
            "slide_count": len(pres.slides),
            "shape_count": len(slide.shapes),
            "content": text_content
        }
    except Exception as e:
        return {
            "error": f"An error occurred: {str(e)}",
            "presentation_id": presentation_id,
            "slide_index": slide_index
        }

@mcp.tool()
def get_slide_shapes(presentation_id: str, slide_index: int) -> Dict[str, Any]:
    """
    Get all shapes in a slide with their IDs and properties.
    
    Args:
        presentation_id: ID of the presentation
        slide_index: Index of the slide (integer, 0-based)
        
    Returns:
        Dictionary containing all shapes with their properties
    """
    try:
        if presentation_id not in ppt_automation.presentations:
            return {"error": f"Presentation ID not found: {presentation_id}"}
        
        pres = ppt_automation.presentations[presentation_id]["presentation"]
        
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {"error": f"Invalid slide index: {slide_index}. Valid range is 0-{len(pres.slides)-1}"}
        
        slide = pres.slides[slide_index]
        shapes_info = {}
        
        for shape_idx, shape in enumerate(slide.shapes):
            shape_id = str(shape_idx)
            shape_type = "Unknown"
            
            # Determine shape type
            if shape.has_text_frame:
                shape_type = "TextFrame"
            elif shape.has_table:
                shape_type = "Table"
            elif shape.has_chart:
                shape_type = "Chart"
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                shape_type = "Picture"
            elif shape.is_placeholder:
                shape_type = f"Placeholder ({shape.placeholder_format.type})"
                
            # Get shape name or generate one
            shape_name = getattr(shape, "name", f"Shape {shape_idx}")
            
            # Basic properties all shapes have
            shape_info = {
                "id": shape_id,
                "name": shape_name,
                "type": shape_type,
                "left": shape.left.inches,
                "top": shape.top.inches,
                "width": shape.width.inches,
                "height": shape.height.inches
            }
            
            # Add text content if available
            if shape.has_text_frame:
                shape_info["text"] = shape.text
                
            shapes_info[shape_id] = shape_info
        
        return {
            "slide_index": slide_index,
            "slide_count": len(pres.slides),
            "shape_count": len(slide.shapes),
            "shapes": shapes_info
        }
    except Exception as e:
        return {
            "error": f"An error occurred: {str(e)}",
            "presentation_id": presentation_id,
            "slide_index": slide_index
        }

@mcp.tool()
def update_text(presentation_id: str, slide_index: int, shape_index: int, text: str) -> Dict[str, Any]:
    """
    Update the text content of a shape.
    
    Args:
        presentation_id: ID of the presentation
        slide_index: Index of the slide (0-based)
        shape_index: Index of the shape (0-based)
        text: New text content
        
    Returns:
        Status of the operation
    """
    if presentation_id not in ppt_automation.presentations:
        return {"error": "Presentation ID not found"}
    
    pres = ppt_automation.presentations[presentation_id]["presentation"]
    
    try:
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {"error": f"Invalid slide index: {slide_index}"}
        
        slide = pres.slides[slide_index]
        
        if shape_index < 0 or shape_index >= len(slide.shapes):
            return {"error": f"Invalid shape index: {shape_index}"}
        
        shape = slide.shapes[shape_index]
        
        if shape.has_text_frame:
            shape.text_frame.text = text
            return {"success": True, "message": "Text updated successfully"}
        else:
            return {"success": False, "message": "Shape does not contain editable text"}
    except Exception as e:
        return {"success": False, "error": f"Error updating text: {str(e)}"}

@mcp.tool()
def update_shape_by_id(presentation_id: str, slide_index: int, shape_id: str, 
                       text: str = None, left: float = None, top: float = None, 
                       width: float = None, height: float = None) -> Dict[str, Any]:
    """
    Update a shape by its ID with new properties.
    
    Args:
        presentation_id: ID of the presentation
        slide_index: Index of the slide (0-based)
        shape_id: ID of the shape (string)
        text: New text content (if applicable)
        left: New left position in inches (if changing)
        top: New top position in inches (if changing)
        width: New width in inches (if changing)
        height: New height in inches (if changing)
        
    Returns:
        Status of the operation
    """
    if presentation_id not in ppt_automation.presentations:
        return {"error": "Presentation ID not found"}
    
    pres = ppt_automation.presentations[presentation_id]["presentation"]
    
    try:
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {"error": f"Invalid slide index: {slide_index}"}
        
        slide = pres.slides[slide_index]
        
        # Convert shape_id to integer index
        try:
            shape_index = int(shape_id)
            if shape_index < 0 or shape_index >= len(slide.shapes):
                return {"error": f"Invalid shape ID: {shape_id}"}
        except ValueError:
            return {"error": f"Invalid shape ID format: {shape_id}. Must be a numeric string."}
        
        shape = slide.shapes[shape_index]
        updates_made = []
        
        # Update text if provided and shape has text frame
        if text is not None and shape.has_text_frame:
            shape.text_frame.text = text
            updates_made.append("text")
        
        # Update position if provided
        if left is not None:
            shape.left = Inches(left)
            updates_made.append("left position")
        
        if top is not None:
            shape.top = Inches(top)
            updates_made.append("top position")
        
        # Update size if provided
        if width is not None:
            shape.width = Inches(width)
            updates_made.append("width")
        
        if height is not None:
            shape.height = Inches(height)
            updates_made.append("height")
        
        if not updates_made:
            return {"success": False, "message": "No updates were specified or applicable to this shape type"}
        
        return {
            "success": True, 
            "message": f"Shape updated successfully. Updated: {', '.join(updates_made)}",
            "shape_id": shape_id
        }
    except Exception as e:
        return {"success": False, "error": f"Error updating shape: {str(e)}"}

@mcp.tool()
def save_presentation(presentation_id: str, path: str = None) -> Dict[str, Any]:
    """
    Save a presentation to disk.
    
    Args:
        presentation_id: ID of the presentation
        path: Optional path to save the file (if None, save to current location)
        
    Returns:
        Status of the operation
    """
    if presentation_id not in ppt_automation.presentations:
        return {"error": "Presentation ID not found"}
    
    pres_info = ppt_automation.presentations[presentation_id]
    pres = pres_info["presentation"]
    
    try:
        save_path = path if path else pres_info["path"]
        
        # If this is a new presentation without a path, we need a path
        if not save_path:
            return {"error": "Save path must be specified for new presentations"}
        
        pres.save(save_path)
        
        # Update the path in our records
        pres_info["path"] = save_path
        
        return {
            "success": True, 
            "path": save_path
        }
    except Exception as e:
        return {"success": False, "error": str(e)}

@mcp.tool()
def close_presentation(presentation_id: str) -> Dict[str, Any]:
    """
    Close a presentation.
    
    Args:
        presentation_id: ID of the presentation
        
    Returns:
        Status of the operation
    """
    if presentation_id not in ppt_automation.presentations:
        return {"error": "Presentation ID not found"}
    
    try:
        # With python-pptx, we just remove it from our tracking
        del ppt_automation.presentations[presentation_id]
        return {"success": True}
    except Exception as e:
        return {"success": False, "error": str(e)}

@mcp.tool()
def create_presentation() -> Dict[str, Any]:
    """
    Create a new PowerPoint presentation.
    
    Returns:
        Dictionary containing new presentation ID and metadata
    """
    try:
        pres = Presentation()
        pres_id = str(uuid.uuid4())
        
        ppt_automation.presentations[pres_id] = {
            "presentation": pres,
            "path": ""
        }
        
        return {
            "id": pres_id,
            "name": "New Presentation",
            "path": "",
            "slide_count": len(pres.slides)
        }
    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def add_slide(presentation_id: str, layout_index: int = 1) -> Dict[str, Any]:
    """
    Add a new slide to the presentation.
    
    Args:
        presentation_id: ID of the presentation
        layout_index: Slide layout index (default is 1)
            0: Title slide
            1: Title and content
            2: Section header
            3: Two content
            etc...
            
    Returns:
        Information about the new slide
    """
    if presentation_id not in ppt_automation.presentations:
        return {"error": "Presentation ID not found"}
    
    pres = ppt_automation.presentations[presentation_id]["presentation"]
    
    try:
        # Get slide layouts
        if layout_index < 0 or layout_index >= len(pres.slide_layouts):
            layout_index = 1  # Default to title and content if invalid
        
        slide_layout = pres.slide_layouts[layout_index]
        
        # Add new slide
        slide = pres.slides.add_slide(slide_layout)
        slide_index = len(pres.slides) - 1
        
        return {
            "id": str(slide_index),
            "index": slide_index,
            "title": "New Slide",
            "shape_count": len(slide.shapes)
        }
    except Exception as e:
        return {"error": f"Error adding slide: {str(e)}"}

@mcp.tool()
def add_textbox(presentation_id: str, slide_index: int, text: str, 
                left: float = 1, top: float = 1, 
                width: float = 4, height: float = 2) -> Dict[str, Any]:
    """
    Add a text box to a slide and set its text content.
    
    Args:
        presentation_id: ID of the presentation
        slide_index: Index of the slide (0-based)
        text: Text content
        left: Left edge position of the text box (inches)
        top: Top edge position of the text box (inches)
        width: Width of the text box (inches)
        height: Height of the text box (inches)
        
    Returns:
        Operation status and ID of the new shape
    """
    if presentation_id not in ppt_automation.presentations:
        return {"error": "Presentation ID not found"}
    
    pres = ppt_automation.presentations[presentation_id]["presentation"]
    
    try:
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {"error": f"Invalid slide index: {slide_index}"}
        
        slide = pres.slides[slide_index]
        
        # Convert to inches
        left_inches = Inches(left)
        top_inches = Inches(top)
        width_inches = Inches(width)
        height_inches = Inches(height)
        
        # Add text box
        shape = slide.shapes.add_textbox(left_inches, top_inches, width_inches, height_inches)
        
        # Set text content
        text_frame = shape.text_frame
        text_frame.text = text
        
        # Get the new shape's index
        shape_index = len(slide.shapes) - 1
        
        return {
            "success": True,
            "slide_index": slide_index,
            "shape_index": shape_index,
            "message": "Text box added successfully"
        }
    except Exception as e:
        return {"error": f"Error adding text box: {str(e)}"}

@mcp.tool()
def set_slide_title(presentation_id: str, slide_index: int, title: str) -> Dict[str, Any]:
    """
    Set the title text of a slide.
    
    Args:
        presentation_id: ID of the presentation
        slide_index: Index of the slide (0-based)
        title: New title text
        
    Returns:
        Status of the operation
    """
    if presentation_id not in ppt_automation.presentations:
        return {"error": "Presentation ID not found"}
    
    pres = ppt_automation.presentations[presentation_id]["presentation"]
    
    try:
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {"error": f"Invalid slide index: {slide_index}"}
        
        slide = pres.slides[slide_index]
        
        # Find title placeholder
        title_found = False
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.type == 1:  # Title placeholder
                shape.text_frame.text = title
                title_found = True
                break
        
        if not title_found:
            # If no title placeholder found, add a text box as title
            left = Inches(1)
            top = Inches(0.5)
            width = Inches(8)
            height = Inches(1)
            
            shape = slide.shapes.add_textbox(left, top, width, height)
            text_frame = shape.text_frame
            text_frame.text = title
            
            # Set text format as title style
            p = text_frame.paragraphs[0]
            p.font.size = Pt(44)
            p.font.bold = True
        
        return {
            "success": True,
            "message": "Slide title has been set"
        }
    except Exception as e:
        return {"error": f"Error setting slide title: {str(e)}"}

@mcp.tool()
def add_image(presentation_id: str, slide_index: int, image_path: str,
              left: float = 1, top: float = 1, width: float = None, height: float = None) -> Dict[str, Any]:
    """
    Add an image to a slide.
    
    Args:
        presentation_id: ID of the presentation
        slide_index: Index of the slide (0-based)
        image_path: Path to the image file
        left: Left position in inches
        top: Top position in inches
        width: Width in inches (optional, maintains aspect ratio if only height is specified)
        height: Height in inches (optional, maintains aspect ratio if only width is specified)
        
    Returns:
        Status of the operation
    """
    if presentation_id not in ppt_automation.presentations:
        return {"error": "Presentation ID not found"}
    
    if not os.path.exists(image_path):
        return {"error": f"Image file not found: {image_path}"}
    
    pres = ppt_automation.presentations[presentation_id]["presentation"]
    
    try:
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {"error": f"Invalid slide index: {slide_index}"}
        
        slide = pres.slides[slide_index]
        
        # Convert to inches
        left_inches = Inches(left)
        top_inches = Inches(top)
        
        # Add image
        if width is not None and height is not None:
            width_inches = Inches(width)
            height_inches = Inches(height)
            shape = slide.shapes.add_picture(image_path, left_inches, top_inches, width_inches, height_inches)
        elif width is not None:
            width_inches = Inches(width)
            shape = slide.shapes.add_picture(image_path, left_inches, top_inches, width=width_inches)
        elif height is not None:
            height_inches = Inches(height)
            shape = slide.shapes.add_picture(image_path, left_inches, top_inches, height=height_inches)
        else:
            shape = slide.shapes.add_picture(image_path, left_inches, top_inches)
        
        # Get the new shape's index
        shape_index = len(slide.shapes) - 1
        
        return {
            "success": True,
            "slide_index": slide_index,
            "shape_index": shape_index,
            "message": "Image added successfully"
        }
    except Exception as e:
        return {"error": f"Error adding image: {str(e)}"}

@mcp.tool()
def add_table(presentation_id: str, slide_index: int, rows: int, cols: int,
              left: float = 1, top: float = 1, width: float = 8, height: float = 4) -> Dict[str, Any]:
    """
    Add a table to a slide.
    
    Args:
        presentation_id: ID of the presentation
        slide_index: Index of the slide (0-based)
        rows: Number of rows
        cols: Number of columns
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        
    Returns:
        Status of the operation
    """
    if presentation_id not in ppt_automation.presentations:
        return {"error": "Presentation ID not found"}
    
    pres = ppt_automation.presentations[presentation_id]["presentation"]
    
    try:
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {"error": f"Invalid slide index: {slide_index}"}
        
        slide = pres.slides[slide_index]
        
        # Convert to inches
        left_inches = Inches(left)
        top_inches = Inches(top)
        width_inches = Inches(width)
        height_inches = Inches(height)
        
        # Add table
        shape = slide.shapes.add_table(rows, cols, left_inches, top_inches, width_inches, height_inches)
        
        # Get the new shape's index
        shape_index = len(slide.shapes) - 1
        
        return {
            "success": True,
            "slide_index": slide_index,
            "shape_index": shape_index,
            "message": "Table added successfully"
        }
    except Exception as e:
        return {"error": f"Error adding table: {str(e)}"}

@mcp.tool()
def update_table_cell(presentation_id: str, slide_index: int, shape_index: int, 
                      row: int, col: int, text: str) -> Dict[str, Any]:
    """
    Update the text in a table cell.
    
    Args:
        presentation_id: ID of the presentation
        slide_index: Index of the slide (0-based)
        shape_index: Index of the table shape (0-based)
        row: Row index (0-based)
        col: Column index (0-based)
        text: New text content
        
    Returns:
        Status of the operation
    """
    if presentation_id not in ppt_automation.presentations:
        return {"error": "Presentation ID not found"}
    
    pres = ppt_automation.presentations[presentation_id]["presentation"]
    
    try:
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {"error": f"Invalid slide index: {slide_index}"}
        
        slide = pres.slides[slide_index]
        
        if shape_index < 0 or shape_index >= len(slide.shapes):
            return {"error": f"Invalid shape index: {shape_index}"}
        
        shape = slide.shapes[shape_index]
        
        if not hasattr(shape, "table"):
            return {"error": "Shape is not a table"}
        
        table = shape.table
        
        if row < 0 or row >= len(table.rows):
            return {"error": f"Invalid row index: {row}"}
        
        if col < 0 or col >= len(table.columns):
            return {"error": f"Invalid column index: {col}"}
        
        cell = table.cell(row, col)
        cell.text = text
        
        return {
            "success": True,
            "message": "Table cell updated successfully"
        }
    except Exception as e:
        return {"error": f"Error updating table cell: {str(e)}"}

@mcp.tool()
def add_chart(presentation_id: str, slide_index: int, chart_type: str,
              categories: List[str], series_names: List[str], series_values: List[List[float]],
              left: float = 1, top: float = 1, width: float = 8, height: float = 4,
              has_legend: bool = True) -> Dict[str, Any]:
    """
    Add a chart to a slide.
    
    Args:
        presentation_id: ID of the presentation
        slide_index: Index of the slide (0-based)
        chart_type: Type of chart ('COLUMN', 'LINE', 'PIE', 'BAR')
        categories: List of category names
        series_names: List of series names
        series_values: List of lists containing values for each series
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        has_legend: Whether to display the legend
        
    Returns:
        Status of the operation
    """
    if presentation_id not in ppt_automation.presentations:
        return {"error": "Presentation ID not found"}
    
    pres = ppt_automation.presentations[presentation_id]["presentation"]
    
    try:
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {"error": f"Invalid slide index: {slide_index}"}
        
        slide = pres.slides[slide_index]
        
        # Convert to inches
        left_inches = Inches(left)
        top_inches = Inches(top)
        width_inches = Inches(width)
        height_inches = Inches(height)
        
        # Map chart type string to ChartType enum
        chart_type_map = {
            'COLUMN': 1,  # XL_CHART_TYPE.COLUMN_CLUSTERED
            'LINE': 4,    # XL_CHART_TYPE.LINE
            'PIE': 5,     # XL_CHART_TYPE.PIE
            'BAR': 57     # XL_CHART_TYPE.BAR_CLUSTERED
        }
        
        chart_type_enum = chart_type_map.get(chart_type.upper(), 1)  # Default to column chart
        
        # Add chart
        chart_data = pres.charts.add_chart(chart_type_enum, left_inches, top_inches, 
                                          width_inches, height_inches).chart_data
        
        # Add categories
        chart_data.categories = categories
        
        # Add series
        for i, (name, values) in enumerate(zip(series_names, series_values)):
            series = chart_data.add_series(name)
            series.values = values
        
        # Get the new shape's index
        shape_index = len(slide.shapes) - 1
        
        return {
            "success": True,
            "slide_index": slide_index,
            "shape_index": shape_index,
            "message": "Chart added successfully"
        }
    except Exception as e:
        return {"error": f"Error adding chart: {str(e)}"}

if __name__ == "__main__":
    mcp.run(transport="stdio")
