from mcp.server.fastmcp import FastMCP
import os
import uuid
from typing import Dict, List, Optional, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
import tempfile

mcp = FastMCP("ppts")

USER_AGENT = "ppts-app/1.0"

class PPTAutomation:
    def __init__(self):
        self.active_presentation = None
        self.presentation_path = None
        
    def initialize(self):
        """Initialize the PowerPoint automation - no app instance needed with python-pptx"""
        return True
                
    def get_active_presentation(self):
        """Get information about the currently active presentation"""
        if self.active_presentation is None:
            return None
        
        return {
            "name": os.path.basename(self.presentation_path) if self.presentation_path else "Untitled",
            "path": self.presentation_path,
            "slide_count": len(self.active_presentation.slides)
        }

# Create a global instance of our automation class
ppt_automation = PPTAutomation()

@mcp.tool()
def initialize_powerpoint() -> bool:
    """Initialize connection to PowerPoint."""
    return ppt_automation.initialize()

@mcp.tool()
def get_presentation() -> Dict[str, Any]:
    """Get information about the currently active presentation."""
    presentation_info = ppt_automation.get_active_presentation()
    if presentation_info is None:
        return {"error": "No active presentation. Please open or create a presentation first."}
    return presentation_info

@mcp.tool()
def open_presentation(file_path: str) -> Dict[str, Any]:
    """
    Open a PowerPoint presentation from the specified path. 
    
    Args:
        file_path: Full path to the PowerPoint file (.pptx)
        
    Returns:
        Dictionary with presentation metadata
    """
    if not os.path.exists(file_path):
        return {"error": f"File not found: {file_path}"}
    
    try:
        ppt_automation.active_presentation = Presentation(file_path)
        ppt_automation.presentation_path = file_path
        
        return {
            "name": os.path.basename(file_path),
            "path": file_path,
            "slide_count": len(ppt_automation.active_presentation.slides)
        }
    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def get_slides() -> List[Dict[str, Any]]:
    """
    Get a list of all slides in the active presentation.
    
    Returns:
        List of slide metadata
    """
    if ppt_automation.active_presentation is None:
        return {"error": "No active presentation. Please open or create a presentation first."}
    
    pres = ppt_automation.active_presentation
    slides = []
    
    try:
        for i, slide in enumerate(pres.slides):
            slide_id = str(i)
            
            # Try to get slide title
            title = "Untitled Slide"
            for shape in slide.shapes:
                if shape.is_placeholder and shape.placeholder_format.type == 1:  # Title placeholder
                    if shape.has_text_frame:
                        title = shape.text
                        break
            
            slides.append({
                "id": slide_id,
                "index": i,
                "title": title,
                "shape_count": len(slide.shapes)
            })
        
        return slides
    except Exception as e:
        return {"error": f"Error getting slides: {str(e)}"}

@mcp.tool()
def get_slide_text(slide_index: int) -> Dict[str, Any]:
    """
    Get all text content in a slide. ALWAYS check if there are groupped shapes first.
    
    Args:
        slide_index: Index of the slide (integer, 0-based)
        
    Returns:
        Dictionary containing text content organized by shape
    """
    try:
        if ppt_automation.active_presentation is None:
            return {"error": "No active presentation. Please open or create a presentation first."}
        
        pres = ppt_automation.active_presentation
        
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {"error": f"Invalid slide index: {slide_index}. Valid range is 0-{len(pres.slides)-1}"}
        
        slide = pres.slides[slide_index]
        text_content = {}
        has_grouped_shapes = False
        
        for shape_idx, shape in enumerate(slide.shapes):
            shape_id = str(shape_idx)
            
            # Check if this is a grouped shape
            if hasattr(shape, 'shape_type') and shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                has_grouped_shapes = True
            
            if shape.has_text_frame:
                text = shape.text
                shape_name = f"Shape {shape_idx}"
                
                text_content[shape_id] = {
                    "shape_name": shape_name,
                    "text": text
                }
        
        return {
            "slide_index": slide_index,
            "slide_count": len(pres.slides),
            "shape_count": len(slide.shapes),
            "has_grouped_shapes": has_grouped_shapes,
            "content": text_content
        }
    except Exception as e:
        return {
            "error": f"An error occurred: {str(e)}",
            "slide_index": slide_index
        }

@mcp.tool()
def get_slide_shapes(slide_index: int) -> Dict[str, Any]:
    """
    ALWAYS RUN IT FIRST!! Get all shapes in a slide with their IDs and properties. If shapes are groupped, remember to ungroup them.
    
    Args:
        slide_index: Index of the slide (integer, 0-based)
        
    Returns:
        Dictionary containing all shapes with their properties
    """
    try:
        if ppt_automation.active_presentation is None:
            return {"error": "No active presentation. Please open or create a presentation first."}
        
        pres = ppt_automation.active_presentation
        
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {"error": f"Invalid slide index: {slide_index}. Valid range is 0-{len(pres.slides)-1}"}
        
        slide = pres.slides[slide_index]
        shapes_info = {}
        
        for shape_idx, shape in enumerate(slide.shapes):
            shape_id = str(shape_idx)
            shape_type = "Unknown"
            
            # Determine shape type
            if shape.has_text_frame:
                shape_type = "TextFrame"
            elif shape.has_table:
                shape_type = "Table"
            elif shape.has_chart:
                shape_type = "Chart"
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                shape_type = "Picture"
            elif shape.is_placeholder:
                shape_type = f"Placeholder ({shape.placeholder_format.type})"
                
            # Get shape name or generate one
            shape_name = getattr(shape, "name", f"Shape {shape_idx}")
            
            # Basic properties all shapes have
            shape_info = {
                "id": shape_id,
                "name": shape_name,
                "type": shape_type,
                "left": shape.left.inches,
                "top": shape.top.inches,
                "width": shape.width.inches,
                "height": shape.height.inches
            }
            
            # Add text content if available
            if shape.has_text_frame:
                shape_info["text"] = shape.text
                
            shapes_info[shape_id] = shape_info
        
        return {
            "slide_index": slide_index,
            "slide_count": len(pres.slides),
            "shape_count": len(slide.shapes),
            "shapes": shapes_info
        }
    except Exception as e:
        return {
            "error": f"An error occurred: {str(e)}",
            "slide_index": slide_index
        }

@mcp.tool()
def update_text(slide_index: int, shape_index: int, text: str,
                font_name: str = None, font_size: int = None,
                bold: bool = None, italic: bool = None,
                preserve_existing: bool = True) -> Dict[str, Any]:
    """
    Update the text content of a shape with optional formatting control.
    
    Args:
        slide_index: Index of the slide (0-based)
        shape_index: Index of the shape (0-based)
        text: New text content
        font_name: Font name (e.g., 'Arial', 'Calibri') - optional
        font_size: Font size in points - optional
        bold: Whether text should be bold - optional
        italic: Whether text should be italic - optional
        preserve_existing: Whether to preserve existing formatting when new formatting is not specified (default: True)
        
    Returns:
        Status of the operation
    """
    if ppt_automation.active_presentation is None:
        return {"error": "No active presentation. Please open or create a presentation first."}
    
    pres = ppt_automation.active_presentation
    
    try:
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {"error": f"Invalid slide index: {slide_index}"}
        
        slide = pres.slides[slide_index]
        
        if shape_index < 0 or shape_index >= len(slide.shapes):
            return {"error": f"Invalid shape index: {shape_index}"}
        
        shape = slide.shapes[shape_index]
        
        if shape.has_text_frame:
            text_frame = shape.text_frame
            
            # Store original formatting if preserving existing
            original_font = None
            original_alignment = None
            if preserve_existing and text_frame.paragraphs and text_frame.paragraphs[0].runs:
                first_run = text_frame.paragraphs[0].runs[0]
                original_alignment = text_frame.paragraphs[0].alignment
                
                # Safely handle color property including theme colors
                color_rgb = None
                theme_color = None
                try:
                    color_obj = first_run.font.color
                    if hasattr(color_obj, 'theme_color') and color_obj.theme_color is not None:
                        # Only save valid theme colors (not NOT_THEME_COLOR)
                        if color_obj.theme_color != 0:  # 0 = NOT_THEME_COLOR
                            theme_color = color_obj.theme_color
                    elif hasattr(color_obj, 'rgb') and color_obj.rgb is not None:
                        color_rgb = color_obj.rgb
                except (AttributeError, TypeError):
                    # Color might be theme-based or None, skip it
                    pass
                
                original_font = {
                    'name': first_run.font.name,
                    'size': first_run.font.size,
                    'bold': first_run.font.bold,
                    'italic': first_run.font.italic,
                    'color': color_rgb,
                    'theme_color': theme_color
                }
            
            # Clear existing content
            text_frame.clear()
            
            # Add new paragraph with the text
            paragraph = text_frame.paragraphs[0]
            
            # Restore alignment if preserved
            if preserve_existing and original_alignment is not None:
                paragraph.alignment = original_alignment
            run = paragraph.add_run()
            run.text = text
            
            # Apply formatting - new parameters take precedence over preserved formatting
            if font_name:
                run.font.name = font_name
            elif preserve_existing and original_font and original_font['name']:
                run.font.name = original_font['name']
            
            if font_size:
                run.font.size = Pt(font_size)
            elif preserve_existing and original_font and original_font['size']:
                run.font.size = original_font['size']
            
            if bold is not None:
                run.font.bold = bold
            elif preserve_existing and original_font and original_font['bold'] is not None:
                run.font.bold = original_font['bold']
            
            if italic is not None:
                run.font.italic = italic
            elif preserve_existing and original_font and original_font['italic'] is not None:
                run.font.italic = original_font['italic']
            
            # Preserve color if it existed (theme color takes priority)
            if preserve_existing and original_font:
                try:
                    if original_font['theme_color'] is not None and original_font['theme_color'] != 0:
                        # Only restore valid theme colors (not NOT_THEME_COLOR)
                        run.font.color.theme_color = original_font['theme_color']
                    elif original_font['color']:
                        run.font.color.rgb = original_font['color']
                except (AttributeError, TypeError):
                    # Skip color if it can't be applied
                    pass
            
            # Build success message
            formatting_applied = []
            if font_name:
                formatting_applied.append(f"font: {font_name}")
            if font_size:
                formatting_applied.append(f"size: {font_size}pt")
            if bold is not None:
                formatting_applied.append(f"bold: {bold}")
            if italic is not None:
                formatting_applied.append(f"italic: {italic}")
            
            message = "Text updated successfully"
            if formatting_applied:
                message += f" with formatting: {', '.join(formatting_applied)}"
            elif preserve_existing:
                message += " with preserved formatting"
            
            return {"success": True, "message": message}
        else:
            return {"success": False, "message": "Shape does not contain editable text"}
    except Exception as e:
        return {"success": False, "error": f"Error updating text: {str(e)}"}

@mcp.tool()
def update_shape_by_id(slide_index: int, shape_id: str, 
                      text: str = None, left: float = None, top: float = None, 
                      width: float = None, height: float = None) -> Dict[str, Any]:
    """
    Update a shape by its ID with new properties.
    
    Args:
        slide_index: Index of the slide (0-based)
        shape_id: ID of the shape (string)
        text: New text content (if applicable)
        left: New left position in inches (if changing)
        top: New top position in inches (if changing)
        width: New width in inches (if changing)
        height: New height in inches (if changing)
        
    Returns:
        Status of the operation
    """
    if ppt_automation.active_presentation is None:
        return {"error": "No active presentation. Please open or create a presentation first."}
    
    pres = ppt_automation.active_presentation
    
    try:
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {"error": f"Invalid slide index: {slide_index}"}
        
        slide = pres.slides[slide_index]
        
        # Convert shape_id to integer index
        try:
            shape_index = int(shape_id)
            if shape_index < 0 or shape_index >= len(slide.shapes):
                return {"error": f"Invalid shape ID: {shape_id}"}
        except ValueError:
            return {"error": f"Invalid shape ID format: {shape_id}. Must be a numeric string."}
        
        shape = slide.shapes[shape_index]
        updates_made = []
        
        # Update text if provided and shape has text frame
        if text is not None and shape.has_text_frame:
            shape.text_frame.text = text
            updates_made.append("text")
        
        # Update position if provided
        if left is not None:
            shape.left = Inches(left)
            updates_made.append("left position")
        
        if top is not None:
            shape.top = Inches(top)
            updates_made.append("top position")
        
        # Update size if provided
        if width is not None:
            shape.width = Inches(width)
            updates_made.append("width")
        
        if height is not None:
            shape.height = Inches(height)
            updates_made.append("height")
        
        if not updates_made:
            return {"success": False, "message": "No updates were specified or applicable to this shape type"}
        
        return {
            "success": True, 
            "message": f"Shape updated successfully. Updated: {', '.join(updates_made)}",
            "shape_id": shape_id
        }
    except Exception as e:
        return {"success": False, "error": f"Error updating shape: {str(e)}"}

@mcp.tool()
def save_presentation(path: str = None) -> Dict[str, Any]:
    """
    Save the active presentation to disk.
    
    Args:
        path: Optional path to save the file (if None, save to current location)
        
    Returns:
        Status of the operation
    """
    if ppt_automation.active_presentation is None:
        return {"error": "No active presentation. Please open or create a presentation first."}
    
    try:
        save_path = path if path else ppt_automation.presentation_path
        
        # If this is a new presentation without a path, we need a path
        if not save_path:
            return {"error": "Save path must be specified for new presentations"}
        
        ppt_automation.active_presentation.save(save_path)
        
        # Update the path in our records
        ppt_automation.presentation_path = save_path
        
        return {
            "success": True, 
            "path": save_path
        }
    except Exception as e:
        return {"success": False, "error": str(e)}

@mcp.tool()
def close_presentation() -> Dict[str, Any]:
    """
    Close the active presentation.
    
    Returns:
        Status of the operation
    """
    if ppt_automation.active_presentation is None:
        return {"error": "No active presentation. Please open or create a presentation first."}
    
    try:
        # With python-pptx, we just remove it from our tracking
        ppt_automation.active_presentation = None
        ppt_automation.presentation_path = None
        return {"success": True}
    except Exception as e:
        return {"success": False, "error": str(e)}

@mcp.tool()
def create_presentation(template: str = None) -> Dict[str, Any]:
    """
    Create a new PowerPoint presentation.
    Args:
        template: Optional path to a template file (.pptx)
    Returns:
        Dictionary containing new presentation metadata
    """
    try:
        ppt_automation.active_presentation = Presentation(pptx=template) if template else Presentation()
        ppt_automation.presentation_path = ""
        
        return {
            "name": "New Presentation",
            "path": "",
            "slide_count": len(ppt_automation.active_presentation.slides)
        }
    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def add_slide(layout_index: int = 1) -> Dict[str, Any]:
    """
    Add a new slide to the presentation.
    
    Args:
        layout_index: Slide layout index (default is 1)
        get_slide_layouts() to see available layouts.
            
    Returns:
        Information about the new slide
    """
    if ppt_automation.active_presentation is None:
        return {"error": "No active presentation. Please open or create a presentation first."}
    
    pres = ppt_automation.active_presentation
    
    try:
        # Get slide layouts
        if layout_index < 0 or layout_index >= len(pres.slide_layouts):
            layout_index = 1  # Default to title and content if invalid
        
        slide_layout = pres.slide_layouts[layout_index]
        
        # Add new slide
        slide = pres.slides.add_slide(slide_layout)
        slide_index = len(pres.slides) - 1
        
        return {
            "id": str(slide_index),
            "index": slide_index,
            "title": "New Slide",
            "shape_count": len(slide.shapes)
        }
    except Exception as e:
        return {"error": f"Error adding slide: {str(e)}"}

@mcp.tool()
def add_textbox(slide_index: int, text: str, 
                left: float = 1, top: float = 1, 
                width: float = 4, height: float = 2) -> Dict[str, Any]:
    """
    Add a text box to a slide and set its text content.
    
    Args:
        slide_index: Index of the slide (0-based)
        text: Text content
        left: Left edge position of the text box (inches)
        top: Top edge position of the text box (inches)
        width: Width of the text box (inches)
        height: Height of the text box (inches)
        
    Returns:
        Operation status and ID of the new shape
    """
    if ppt_automation.active_presentation is None:
        return {"error": "No active presentation. Please open or create a presentation first."}
    
    pres = ppt_automation.active_presentation
    
    try:
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {"error": f"Invalid slide index: {slide_index}"}
        
        slide = pres.slides[slide_index]
        
        # Convert to inches
        left_inches = Inches(left)
        top_inches = Inches(top)
        width_inches = Inches(width)
        height_inches = Inches(height)
        
        # Add text box
        shape = slide.shapes.add_textbox(left_inches, top_inches, width_inches, height_inches)
        
        # Set text content
        text_frame = shape.text_frame
        text_frame.text = text
        
        # Get the new shape's index
        shape_index = len(slide.shapes) - 1
        
        return {
            "success": True,
            "slide_index": slide_index,
            "shape_index": shape_index,
            "message": "Text box added successfully"
        }
    except Exception as e:
        return {"error": f"Error adding text box: {str(e)}"}

@mcp.tool()
def set_slide_title(slide_index: int, title: str) -> Dict[str, Any]:
    """
    Set the title text of a slide.
    
    Args:
        slide_index: Index of the slide (0-based)
        title: New title text
        
    Returns:
        Status of the operation
    """
    if ppt_automation.active_presentation is None:
        return {"error": "No active presentation. Please open or create a presentation first."}
    
    pres = ppt_automation.active_presentation
    
    try:
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {"error": f"Invalid slide index: {slide_index}"}
        
        slide = pres.slides[slide_index]
        
        # Find title placeholder
        title_found = False
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.type == 1:  # Title placeholder
                shape.text_frame.text = title
                title_found = True
                break
        
        if not title_found:
            # If no title placeholder found, add a text box as title
            left = Inches(1)
            top = Inches(0.5)
            width = Inches(8)
            height = Inches(1)
            
            shape = slide.shapes.add_textbox(left, top, width, height)
            text_frame = shape.text_frame
            text_frame.text = title
            
            # Set text format as title style
            p = text_frame.paragraphs[0]
            p.font.size = Pt(44)
            p.font.bold = True
        
        return {
            "success": True,
            "message": "Slide title has been set"
        }
    except Exception as e:
        return {"error": f"Error setting slide title: {str(e)}"}

@mcp.tool()
def add_image(slide_index: int, image_path: str,
              left: float = 1, top: float = 1, width: float = None, height: float = None) -> Dict[str, Any]:
    """
    Add an image to a slide.
    
    Args:
        slide_index: Index of the slide (0-based)
        image_path: Path to the image file
        left: Left position in inches
        top: Top position in inches
        width: Width in inches (optional, maintains aspect ratio if only height is specified)
        height: Height in inches (optional, maintains aspect ratio if only width is specified)
        
    Returns:
        Status of the operation
    """
    if ppt_automation.active_presentation is None:
        return {"error": "No active presentation. Please open or create a presentation first."}
    
    if not os.path.exists(image_path):
        return {"error": f"Image file not found: {image_path}"}
    
    pres = ppt_automation.active_presentation
    
    try:
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {"error": f"Invalid slide index: {slide_index}"}
        
        slide = pres.slides[slide_index]
        
        # Convert to inches
        left_inches = Inches(left)
        top_inches = Inches(top)
        
        # Add image
        if width is not None and height is not None:
            width_inches = Inches(width)
            height_inches = Inches(height)
            shape = slide.shapes.add_picture(image_path, left_inches, top_inches, width_inches, height_inches)
        elif width is not None:
            width_inches = Inches(width)
            shape = slide.shapes.add_picture(image_path, left_inches, top_inches, width=width_inches)
        elif height is not None:
            height_inches = Inches(height)
            shape = slide.shapes.add_picture(image_path, left_inches, top_inches, height=height_inches)
        else:
            shape = slide.shapes.add_picture(image_path, left_inches, top_inches)
        
        # Get the new shape's index
        shape_index = len(slide.shapes) - 1
        
        return {
            "success": True,
            "slide_index": slide_index,
            "shape_index": shape_index,
            "message": "Image added successfully"
        }
    except Exception as e:
        return {"error": f"Error adding image: {str(e)}"}

@mcp.tool()
def add_table(slide_index: int, rows: int, cols: int,
              left: float = 1, top: float = 1, width: float = 8, height: float = 4) -> Dict[str, Any]:
    """
    Add a table to a slide.
    
    Args:
        slide_index: Index of the slide (0-based)
        rows: Number of rows
        cols: Number of columns
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        
    Returns:
        Status of the operation
    """
    if ppt_automation.active_presentation is None:
        return {"error": "No active presentation. Please open or create a presentation first."}
    
    pres = ppt_automation.active_presentation
    
    try:
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {"error": f"Invalid slide index: {slide_index}"}
        
        slide = pres.slides[slide_index]
        
        # Convert to inches
        left_inches = Inches(left)
        top_inches = Inches(top)
        width_inches = Inches(width)
        height_inches = Inches(height)
        
        # Add table
        shape = slide.shapes.add_table(rows, cols, left_inches, top_inches, width_inches, height_inches)
        
        # Get the new shape's index
        shape_index = len(slide.shapes) - 1
        
        return {
            "success": True,
            "slide_index": slide_index,
            "shape_index": shape_index,
            "message": "Table added successfully"
        }
    except Exception as e:
        return {"error": f"Error adding table: {str(e)}"}

@mcp.tool()
def update_table_cell(slide_index: int, shape_index: int, 
                      row: int, col: int, text: str) -> Dict[str, Any]:
    """
    Update the text in a table cell.
    
    Args:
        slide_index: Index of the slide (0-based)
        shape_index: Index of the table shape (0-based)
        row: Row index (0-based)
        col: Column index (0-based)
        text: New text content
        
    Returns:
        Status of the operation
    """
    if ppt_automation.active_presentation is None:
        return {"error": "No active presentation. Please open or create a presentation first."}
    
    pres = ppt_automation.active_presentation
    print(f"Active presentation: {ppt_automation.active_presentation}")
    try:
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {"error": f"Invalid slide index: {slide_index}"}
        
        slide = pres.slides[slide_index]
        
        if shape_index < 0 or shape_index >= len(slide.shapes):
            return {"error": f"Invalid shape index: {shape_index}"}
        
        shape = slide.shapes[shape_index]
        
        if not hasattr(shape, "table"):
            return {"error": "Shape is not a table"}
        
        table = shape.table
        if row < 0 or row >= len(table.rows):
            return {"error": f"Invalid row index: {row}"}
        
        if col < 0 or col >= len(table.columns):
            return {"error": f"Invalid column index: {col}"}
        
        cell = table.cell(row, col)
        cell.text = text
        print(f"Updated cell ({row}, {col}) with text: {text}")
        return {
            "success": True,
            "message": "Table cell updated successfully"
        }
    except Exception as e:
        return {"error": f"Error updating table cell: {str(e)}"}

@mcp.tool()
def get_table_content(slide_index: int, shape_index: int) -> Dict[str, Any]:
    """
    Retrieve the content of a table in a slide.
    
    Args:
        slide_index: Index of the slide (0-based)
        shape_index: Index of the table shape (0-based)
        
    Returns:
        Dictionary containing table data with rows and columns
    """
    if ppt_automation.active_presentation is None:
        return {"error": "No active presentation. Please open or create a presentation first."}
    
    pres = ppt_automation.active_presentation
    
    try:
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {"error": f"Invalid slide index: {slide_index}"}
        
        slide = pres.slides[slide_index]
        
        if shape_index < 0 or shape_index >= len(slide.shapes):
            return {"error": f"Invalid shape index: {shape_index}"}
        
        shape = slide.shapes[shape_index]
        
        if not hasattr(shape, "table"):
            return {"error": "Shape is not a table"}
        
        table = shape.table
        rows_count = len(table.rows)
        cols_count = len(table.columns)
        
        # Extract table data
        table_data = []
        for row_idx in range(rows_count):
            row_data = []
            for col_idx in range(cols_count):
                cell = table.cell(row_idx, col_idx)
                row_data.append(cell.text)
            table_data.append(row_data)
        
        return {
            "success": True,
            "rows": rows_count,
            "columns": cols_count,
            "data": table_data
        }
    except Exception as e:
        return {"error": f"Error retrieving table content: {str(e)}"}

@mcp.tool()
def get_slide_layouts() -> Dict[str, Any]:
    """
    Get all available slide layouts in the active presentation.
    
    Returns:
        Dictionary containing all slide layouts with their properties
    """
    if ppt_automation.active_presentation is None:
        return {"error": "No active presentation. Please open or create a presentation first."}
    
    pres = ppt_automation.active_presentation
    
    try:
        layouts = []
        for i, layout in enumerate(pres.slide_layouts):
            # Try to get layout name
            layout_name = layout.name if hasattr(layout, "name") else f"Layout {i}"
            
            # Count placeholders
            placeholder_count = 0
            placeholder_types = []
            for shape in layout.placeholders:
                placeholder_count += 1
                ph_type = shape.placeholder_format.type
                placeholder_types.append(ph_type)
            
            layouts.append({
                "index": i,
                "name": layout_name,
                "placeholder_count": placeholder_count,
                "placeholder_types": placeholder_types
            })
        
        return {
            "success": True,
            "layout_count": len(pres.slide_layouts),
            "layouts": layouts
        }
    except Exception as e:
        return {"error": f"Error getting slide layouts: {str(e)}"}

@mcp.tool()
def delete_slide(slide_index: int) -> Dict[str, Any]:
    """
    Delete a slide from the presentation.
    
    Args:
        slide_index: Index of the slide to delete (0-based)
        
    Returns:
        Status of the operation
    """
    if ppt_automation.active_presentation is None:
        return {"error": "No active presentation. Please open or create a presentation first."}
    
    pres = ppt_automation.active_presentation
    
    try:
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {"error": f"Invalid slide index: {slide_index}. Valid range is 0-{len(pres.slides)-1}"}
        
        # Get the XML slides list
        slides = pres.slides._sldIdLst
        
        # Get the slide ID to remove
        slide_id = slides[slide_index].rId
        
        # Remove the slide from the list
        slides.remove(slides[slide_index])
        
        # Remove the relationship
        pres.part.rels.remove(slide_id)
        
        return {
            "success": True,
            "message": f"Slide at index {slide_index} has been deleted",
            "remaining_slides": len(pres.slides)
        }
    except Exception as e:
        return {"error": f"Error deleting slide: {str(e)}"}

@mcp.tool()
def add_chart(slide_index: int, chart_type: str,
              categories: List[str], series_names: List[str], series_values: List[List[float]],
              left: float = 1, top: float = 1, width: float = 8, height: float = 4,
              has_legend: bool = True) -> Dict[str, Any]:
    """
    Add a chart to a slide.
    
    Args:
        slide_index: Index of the slide (0-based)
        chart_type: Type of chart ('COLUMN', 'LINE', 'PIE', 'BAR')
        categories: List of category names
        series_names: List of series names
        series_values: List of lists containing values for each series
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        has_legend: Whether to display the legend
        
    Returns:
        Status of the operation
    """
    if ppt_automation.active_presentation is None:
        return {"error": "No active presentation. Please open or create a presentation first."}
    
    pres = ppt_automation.active_presentation
    
    try:
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {"error": f"Invalid slide index: {slide_index}"}
        
        slide = pres.slides[slide_index]
        
        # Convert to inches
        left_inches = Inches(left)
        top_inches = Inches(top)
        width_inches = Inches(width)
        height_inches = Inches(height)
        
        # Map chart type string to ChartType enum
        chart_type_map = {
            'COLUMN': 1,  # XL_CHART_TYPE.COLUMN_CLUSTERED
            'LINE': 4,    # XL_CHART_TYPE.LINE
            'PIE': 5,     # XL_CHART_TYPE.PIE
            'BAR': 57     # XL_CHART_TYPE.BAR_CLUSTERED
        }
        
        chart_type_enum = chart_type_map.get(chart_type.upper(), 1)  # Default to column chart
        
        # Create chart data
        from pptx.chart.data import CategoryChartData
        chart_data = CategoryChartData()
        
        # Add categories
        chart_data.categories = categories
        
        # Add series
        for i, (name, values) in enumerate(zip(series_names, series_values)):
            chart_data.add_series(name, values)
            
        # Add chart to slide
        chart = slide.shapes.add_chart(
            chart_type_enum, left_inches, top_inches, width_inches, height_inches, chart_data
        )
        
        # Get the new shape's index
        shape_index = len(slide.shapes) - 1
        
        return {
            "success": True,
            "slide_index": slide_index,
            "shape_index": shape_index,
            "message": "Chart added successfully"
        }
    except Exception as e:
        return {"error": f"Error adding chart: {str(e)}"}

@mcp.tool()
def ungroup_shapes(slide_index: int) -> Dict[str, Any]:
    """
    Ungroup all groups in a slide (only if all groups are simple without nested grpSp elements).
    
    Args:
        slide_index: Index of the slide (0-based)
        
    Returns:
        Status of the operation
    """
    if ppt_automation.active_presentation is None:
        return {"error": "No active presentation. Please open or create a presentation first."}
    
    pres = ppt_automation.active_presentation
    
    try:
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {"error": f"Invalid slide index: {slide_index}"}
        
        slide = pres.slides[slide_index]
        
        # 檢查所有群組是否都是簡單群組，並標記包含文字的群組
        text_groups = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                group_elem = shape._element
                children = list(group_elem)
                
                has_grpSp = any(child.tag.endswith('}grpSp') for child in children)
                if has_grpSp:
                    return {"error": "Slide contains complex nested groups (grpSp). Entire slide skipped."}
                
                # 檢查群組是否包含文字
                has_text = False
                try:
                    for child_shape in shape.shapes:
                        if hasattr(child_shape, 'has_text_frame') and child_shape.has_text_frame:
                            if child_shape.text.strip():
                                has_text = True
                                break
                except:
                    pass
                
                if has_text:
                    text_groups.append(shape)
        
        if not text_groups:
            return {"success": True, "message": "No text-containing groups found in slide"}
        
        # 反覆處理直到沒有包含文字的群組
        total_groups = 0
        total_shapes = 0
        
        while True:
            has_text_group = False
            shapes_to_process = list(slide.shapes)
            
            for shape in shapes_to_process:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    # 檢查這個群組是否包含文字
                    has_text = False
                    try:
                        for child_shape in shape.shapes:
                            if hasattr(child_shape, 'has_text_frame') and child_shape.has_text_frame:
                                if child_shape.text.strip():
                                    has_text = True
                                    break
                    except:
                        pass
                    
                    if not has_text:
                        continue  # 跳過不包含文字的群組
                    
                    has_text_group = True
                    total_groups += 1
                    
                    spTree = slide._element.spTree
                    group_elem = shape._element
                    
                    # 獲取群組的位置和變換資訊
                    grpSpPr = group_elem.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}grpSpPr')
                    
                    if grpSpPr is not None:
                        xfrm = grpSpPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
                        chOff = grpSpPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}chOff')
                        
                        group_x = int(xfrm.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}off').get('x', 0))
                        group_y = int(xfrm.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}off').get('y', 0))
                        
                        child_off_x = int(chOff.get('x', 0)) if chOff is not None else 0
                        child_off_y = int(chOff.get('y', 0)) if chOff is not None else 0
                    else:
                        group_x = group_y = child_off_x = child_off_y = 0
                    
                    # 處理子元素
                    children = list(group_elem)
                    for child in children:
                        if not child.tag.endswith('}sp'):
                            continue
                            
                        spPr = child.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
                        
                        if spPr is not None:
                            off = spPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}off')
                            
                            if off is not None:
                                child_x = int(off.get('x', 0))
                                child_y = int(off.get('y', 0))
                                
                                new_x = group_x + (child_x - child_off_x)
                                new_y = group_y + (child_y - child_off_y)
                                
                                off.set('x', str(new_x))
                                off.set('y', str(new_y))
                        
                        spTree.append(child)
                        total_shapes += 1
                    
                    spTree.remove(group_elem)
                    break
            
            if not has_text_group:
                break
        
        if total_groups == 0:
            return {"success": True, "message": "No text-containing groups found in slide"}
        
        return {
            "success": True,
            "message": f"Ungrouped {total_groups} text-containing groups, extracted {total_shapes} shapes"
        }
    except Exception as e:
        return {"error": f"Error ungrouping shapes: {str(e)}"}

def main():
    mcp.run(transport="stdio")

if __name__ == "__main__":
    main()
